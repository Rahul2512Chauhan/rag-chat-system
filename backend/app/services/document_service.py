# backend/app/services/document_service.py
from typing import List
import os
import pandas as pd
from langchain.schema import Document

# PDF
try:
    from pypdf import PdfReader  # lightweight, good enough for text
except Exception:
    PdfReader = None

# PPTX
try:
    from pptx import Presentation
except Exception:
    Presentation = None

# DOCX
try:
    import docx
except Exception:
    docx = None


def _extract_pdf(file_path: str) -> List[Document]:
    if PdfReader is None:
        raise RuntimeError("pypdf is not installed")
    docs: List[Document] = []
    reader = PdfReader(file_path)
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            docs.append(
                Document(
                    page_content=text,
                    metadata={"source": os.path.basename(file_path), "page": i + 1},
                )
            )
    return docs


def _extract_pptx(file_path: str) -> List[Document]:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed")

    docs: List[Document] = []
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            # Safe check: both attribute existence and truthy
            if getattr(shape, "has_text_frame", False):
                tf = getattr(shape, "text_frame", None)
                if tf and getattr(tf, "text", "").strip():
                    texts.append(tf.text.strip())
        page_text = "\n".join(texts)
        if page_text:
            docs.append(
                Document(
                    page_content=page_text,
                    metadata={"source": os.path.basename(file_path), "slide": i + 1},
                )
            )
    return docs



def _extract_docx(file_path: str) -> List[Document]:
    if docx is None:
        raise RuntimeError("python-docx is not installed")
    d = docx.Document(file_path)
    paragraphs = [p.text for p in d.paragraphs if p.text and p.text.strip()]
    text = "\n".join(paragraphs)
    if text.strip():
        return [
            Document(
                page_content=text,
                metadata={"source": os.path.basename(file_path)},
            )
        ]
    return []


def _extract_txt(file_path: str) -> List[Document]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    if text.strip():
        return [
            Document(
                page_content=text,
                metadata={"source": os.path.basename(file_path)},
            )
        ]
    return []


def _extract_csv_xlsx(file_path: str) -> List[Document]:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".csv":
            df = pd.read_csv(file_path)
        else:
            # openpyxl engine is already in requirements
            df = pd.read_excel(file_path, engine="openpyxl")
    except Exception:
        return []

    # Convert to CSV text for retrieval (simple & effective)
    text = df.to_csv(index=False)
    if text.strip():
        return [
            Document(
                page_content=text,
                metadata={"source": os.path.basename(file_path), "type": ext},
            )
        ]
    return []


def extract_text_from_file(file_path: str) -> str:
    docs = load_documents(file_path)
    return "\n\n".join([d.page_content for d in docs])


def load_documents(file_path: str) -> List[Document]:
    """
    Load file and return a list of langchain.schema.Document objects.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"{file_path} not found")

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext in [".pptx", ".ppt"]:
        return _extract_pptx(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".doc":
        # Legacy .doc not supported by python-docx
        raise RuntimeError(".doc format is not supported; please convert to .docx")
    elif ext == ".txt":
        return _extract_txt(file_path)
    elif ext in [".csv", ".xlsx", ".xls"]:
        return _extract_csv_xlsx(file_path)
    else:
        # Unknown extension, attempt to read as text
        try:
            return _extract_txt(file_path)
        except Exception:
            return []
